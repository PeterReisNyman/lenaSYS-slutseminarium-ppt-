# -*- coding: utf-8 -*-
"""Generate supplemental_seminar.pptx from slides.json
Run:  python docs/sem_ppt_plan/build_ppt.py
"""
import json
import os
import sys
import subprocess

from pptx import Presentation
from pptx.util import Inches, Pt

BASE = os.path.dirname(__file__)
JSON_PATH = os.path.join(BASE, "slides.json")
OUT_PATH = os.path.join(BASE, "supplemental_seminar.pptx")


def _repo_root(start=BASE):
    try:
        return subprocess.check_output([
            "git", "-C", start, "rev-parse", "--show-toplevel"
        ], text=True).strip()
    except subprocess.CalledProcessError:
        return os.getcwd()


def add_title_slide(prs, item):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
    title_shape = slide.shapes.title
    title_shape.text = f"Issue/PR {item['id']} – {item['title']}"
    if title_shape.text_frame.paragraphs:
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()  # wipe default bullets

    meta = [
        f"Commit: {item.get('hash', '')}",
        f"Date:   {item.get('date', '')}",
        f"Δ LOC:  {item.get('stat', '')}",
        "",
    ] + item.get("bullets", [])

    for line in meta:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)


def add_screenshot_slide(prs, item):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Before     →     After"
        if title_shape.text_frame.paragraphs:
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)

    left = Inches(0.2)
    top = Inches(1.2)
    width = Inches(4.8)

    before_path = os.path.join(BASE, item.get("before", ""))
    after_path = os.path.join(BASE, item.get("after", ""))

    if os.path.isfile(before_path):
        slide.shapes.add_picture(before_path, left, top, width=width)
    else:
        ph = slide.shapes.add_textbox(left, top, width, Inches(3))
        ph.text_frame.text = "PLACEHOLDER"

    if os.path.isfile(after_path):
        slide.shapes.add_picture(after_path, left + Inches(5.1), top, width=width)
    else:
        ph = slide.shapes.add_textbox(left + Inches(5.1), top, width, Inches(3))
        ph.text_frame.text = "PLACEHOLDER"


def add_diff_slide(prs, item, repo_root):
    hash_ = item.get("hash")
    if not hash_:
        return
    try:
        diff = subprocess.check_output([
            "git", "-C", repo_root, "show", hash_
        ], text=True, stderr=subprocess.STDOUT)
    except subprocess.CalledProcessError as e:
        diff = e.output
    lines = diff.splitlines()[:20]
    content = "\n".join(lines)

    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
    slide.shapes.title.text = "Key Diff"
    if slide.shapes.title.text_frame.paragraphs:
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = content
    p.font.name = "Courier New"
    p.font.size = Pt(14)


def main():
    try:
        data = json.load(open(JSON_PATH, encoding="utf-8"))
    except FileNotFoundError:
        sys.exit(f"slides.json not found at {JSON_PATH}")

    repo_root = _repo_root()
    prs = Presentation()
    for item in data:
        add_title_slide(prs, item)
        add_screenshot_slide(prs, item)
        add_diff_slide(prs, item, repo_root)
    prs.save(OUT_PATH)
    print(f"PowerPoint saved → {OUT_PATH}")


if __name__ == "__main__":
    main()
